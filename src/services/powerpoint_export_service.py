from io import BytesIO
from pathlib import Path
from typing import Optional

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.slide import Slide
from pptx.table import Table, _Cell
from pptx.util import Inches, Pt

from src.core.logging import logger
from src.models.address import Coordinates
from src.models.market_analysis import MarketAnalysis
from src.services.google_location_service import GoogleLocationService


LOGO_PATH = Path(__file__).parent.parent.parent / "resources" / "logo.png"


class PowerPointStyles:
    LIGHT_BLUE = RGBColor(219, 233, 247)
    BLUE = RGBColor(91, 155, 213)
    DARK_BLUE = RGBColor(0, 112, 192)
    WHITE = RGBColor(255, 255, 255)
    BLACK = RGBColor(0, 0, 0)
    DEFAULT_FONT = "Arial"


class PowerPointExportService:

    def __init__(
        self,
        location_service: Optional[GoogleLocationService] = None,
        satellite_zoom_level: int = 15,
    ) -> None:
        self._location_service = location_service
        self._satellite_zoom_level = satellite_zoom_level

    def export_market_analysis(self, market_analysis: MarketAnalysis) -> bytes:
        logger.info(f"Building PowerPoint presentation for {market_analysis.address}")

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        self._write_title_slide(prs, market_analysis)

        self._write_second_slide(prs, market_analysis)

        self._write_third_slide(prs, market_analysis)

        output = BytesIO()
        prs.save(output)
        output.seek(0)

        logger.info("PowerPoint presentation built successfully with 3 slides")
        return output.getvalue()

    def _write_title_slide(self, prs: Presentation, market_analysis: MarketAnalysis) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = market_analysis.address
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = PowerPointStyles.BLUE

        address_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
        address_frame = address_box.text_frame
        address_frame.text = f"Address: {market_analysis.address} Placer POI: {market_analysis.reference_poi_name}"
        address_para = address_frame.paragraphs[0]
        address_para.alignment = PP_ALIGN.CENTER
        address_para.font.size = Pt(24)

        self._add_logo_to_slide(slide)

    def _write_second_slide(self, prs: Presentation, market_analysis: MarketAnalysis) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = market_analysis.address
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = PowerPointStyles.BLUE

        border_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.85), Inches(9), Inches(0.03))
        border_line.fill.solid()
        border_line.fill.fore_color.rgb = PowerPointStyles.DARK_BLUE
        border_line.line.fill.background()
        border_line.shadow.inherit = False
        border_line.shadow.visible = False

        left_column_x = Inches(0.5)
        left_column_width = Inches(4.0)
        right_column_x = Inches(5.0)
        right_column_width = Inches(4.5)

        key_stats_top = Inches(1.0)
        tam_top = Inches(2.3)
        site_score_top = Inches(3.4)

        satellite_top = Inches(1.0)
        drive_time_top = Inches(3.75)
        budget_top = Inches(4.95)

        # Left column tables
        self._add_key_stats_cell(slide, market_analysis, left_column_x, key_stats_top, left_column_width, Inches(1.0))
        self._add_tam_cell(slide, market_analysis, left_column_x, tam_top, left_column_width, Inches(1.2))
        self._add_site_score_cell(slide, market_analysis, left_column_x, site_score_top, left_column_width, Inches(2.35))

        # Right column elements
        self._add_satellite_cell(slide, market_analysis, right_column_x, satellite_top, right_column_width, Inches(2.75))
        self._add_drive_time_cell(slide, market_analysis, right_column_x, drive_time_top, right_column_width, Inches(1.2))
        self._add_budget_cell(slide, market_analysis, right_column_x, budget_top, right_column_width, Inches(1.45))

        self._add_logo_to_slide(slide)

    def _write_third_slide(self, prs: Presentation, market_analysis: MarketAnalysis) -> None:
        if not self._location_service:
            logger.info("No location service provided, skipping images slide")
            return

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = market_analysis.address
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = PowerPointStyles.BLUE

        border_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.85), Inches(9), Inches(0.03))
        border_line.fill.solid()
        border_line.fill.fore_color.rgb = PowerPointStyles.DARK_BLUE
        border_line.line.fill.background()
        border_line.shadow.inherit = False
        border_line.shadow.visible = False

        coordinates = Coordinates(
            latitude=market_analysis.latitude,
            longitude=market_analysis.longitude,
            formatted_address=market_analysis.address,
        )

        street_view_bytes = self._location_service.download_street_view_image(coordinates)
        satellite_bytes = self._location_service.download_satellite_image(coordinates, self._satellite_zoom_level)

        # 2x2 grid dimensions
        cell_width = Inches(4.5)
        left_col = Inches(0.5)
        right_col = Inches(5.25)
        top_row = Inches(1.0)
        bottom_row = Inches(4.1)

        if satellite_bytes:
            satellite_stream = BytesIO(satellite_bytes)
            slide.shapes.add_picture(satellite_stream, left_col, top_row, width=cell_width, height=Inches(2.75))
        else:
            logger.warning("Failed to download satellite image")

        if street_view_bytes:
            street_view_stream = BytesIO(street_view_bytes)
            slide.shapes.add_picture(street_view_stream, right_col, top_row, width=cell_width, height=Inches(2.75))
        else:
            logger.warning("Failed to download street view image")

        label_box = slide.shapes.add_textbox(left_col, bottom_row, cell_width, Inches(0.3))
        label_frame = label_box.text_frame
        label_frame.text = "Retail Performance"
        label_para = label_frame.paragraphs[0]
        label_para.font.size = Pt(14)
        label_para.font.bold = True

        if market_analysis.retailers or market_analysis.reference_poi_retail:
            display_retailers = []
            if market_analysis.reference_poi_retail:
                display_retailers.append(market_analysis.reference_poi_retail)
            display_retailers.extend(market_analysis.retailers[: 7 if market_analysis.reference_poi_retail else 8])

            rows = len(display_retailers) + 1
            cols = 4

            table = slide.shapes.add_table(rows, cols, left_col, bottom_row + Inches(0.35), cell_width, Inches(2.4)).table
            self._reset_table_styles(table)

            table.cell(0, 0).text = "Retailer"
            table.cell(0, 1).text = "National %"
            table.cell(0, 2).text = "State %"
            table.cell(0, 3).text = "Distance"

            for col_idx in range(cols):
                cell = table.cell(0, col_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = PowerPointStyles.WHITE
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = PowerPointStyles.BLACK
                        run.font.bold = True
                        run.font.size = Pt(8)

            for i, retailer in enumerate(display_retailers, start=1):
                is_poi = market_analysis.reference_poi_retail and i == 1
                table.cell(i, 0).text = f"* {retailer.name}" if is_poi else retailer.name

                if retailer.national_percentile is not None:
                    percentile = retailer.national_percentile
                    table.cell(i, 1).text = f"{percentile * 100:.0f}%"
                    self._apply_color_scale(table.cell(i, 1), percentile)
                else:
                    table.cell(i, 1).text = "N/A"

                if retailer.state_percentile is not None:
                    percentile = retailer.state_percentile
                    table.cell(i, 2).text = f"{percentile * 100:.0f}%"
                    self._apply_color_scale(table.cell(i, 2), percentile)
                else:
                    table.cell(i, 2).text = "N/A"

                if is_poi:
                    table.cell(i, 3).text = "-"
                else:
                    retailer_distance = retailer.distance_miles
                    table.cell(i, 3).text = f"{retailer_distance:.1f}"
                    distance_value = max(0, 1 - (retailer_distance / 3))
                    self._apply_color_scale(table.cell(i, 3), distance_value)

                for col_idx in range(cols):
                    cell = table.cell(i, col_idx)
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(7)
                            if is_poi:
                                run.font.bold = True
        else:
            text_box = slide.shapes.add_textbox(left_col, bottom_row + Inches(0.35), cell_width, Inches(2.4))
            text_frame = text_box.text_frame
            text_frame.text = "No retail data available"
            text_frame.paragraphs[0].font.size = Pt(12)

        if market_analysis.competitors:
            sorted_competitors = sorted(market_analysis.competitors, key=lambda x: x.members_in_market, reverse=True)
            top_competitors = sorted_competitors[:8]
            rows = len(top_competitors) + 2
            cols = 4

            table = slide.shapes.add_table(rows, cols, right_col, bottom_row + Inches(0.35), cell_width, Inches(2.4)).table
            self._reset_table_styles(table)

            table.cell(0, 0).text = "Competitor"
            table.cell(0, 1).text = "Total"
            table.cell(0, 2).text = "% Overlap"
            table.cell(0, 3).text = "In Market"

            for col_idx in range(cols):
                cell = table.cell(0, col_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = PowerPointStyles.WHITE
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = PowerPointStyles.BLACK
                        run.font.bold = True
                        run.font.size = Pt(8)

            for i, competitor in enumerate(top_competitors, start=1):
                table.cell(i, 0).text = competitor.name
                table.cell(i, 1).text = f"{competitor.total_members:,}"
                table.cell(i, 2).text = f"{competitor.overlap_percentage:.0f}%"
                table.cell(i, 3).text = f"{competitor.members_in_market:,}"

                for col_idx in range(cols):
                    cell = table.cell(i, col_idx)
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(7)

            total_row = len(top_competitors) + 1
            total_members = sum([competitor.members_in_market for competitor in top_competitors])
            table.cell(total_row, 0).text = "Total Members"
            table.cell(total_row, 3).text = f"{total_members:,}"

            for col_idx in range(cols):
                cell = table.cell(total_row, col_idx)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                        run.font.size = Pt(7)
        else:
            text_box = slide.shapes.add_textbox(right_col, bottom_row + Inches(0.35), cell_width, Inches(2.4))
            text_frame = text_box.text_frame
            text_frame.text = "No competitors found"
            text_frame.paragraphs[0].font.size = Pt(12)

        self._add_logo_to_slide(slide)

    def _add_key_stats_cell(self, slide: Slide, market_analysis: MarketAnalysis, left: Inches, top: Inches, width: Inches, height: Inches) -> None:
        result_12min = next((r for r in market_analysis.car_parc_results if r.drive_time_minutes == 12), None)
        car_parc_12min = result_12min.car_parc if result_12min else 0

        stats_data = [
            ("Car Counts", f"{market_analysis.key_stats.car_counts:,}" if market_analysis.key_stats.car_counts else "N/A", "N/A"),
            ("Car Parc", f"{car_parc_12min:,}", "N/A"),
            ("Median Income", f"${market_analysis.key_stats.median_income:,}" if market_analysis.key_stats.median_income else "N/A", "N/A"),
            ("Average Age", f"{round(market_analysis.key_stats.median_age)}" if market_analysis.key_stats.median_age else "N/A", "N/A"),
        ]

        rows = len(stats_data) + 1
        cols = 3
        table = slide.shapes.add_table(rows, cols, left, top, width, Inches(1.0)).table
        self._reset_table_styles(table)

        table.columns[0].width = Inches(1.5)
        table.columns[1].width = Inches(1.5)
        table.columns[2].width = Inches(1.5)

        headers = ["Key Stats", self._get_short_address(market_analysis.address), "Median"]
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = PowerPointStyles.BLACK
                    run.font.bold = True
                    run.font.size = Pt(8)

        for i, (label, portsmouth_value, median_value) in enumerate(stats_data, start=1):
            cell = table.cell(i, 0)
            cell.text = label
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = PowerPointStyles.BLACK
                    run.font.bold = True
                    run.font.size = Pt(7)

            cell = table.cell(i, 1)
            cell.text = portsmouth_value
            cell.fill.solid()
            cell.fill.fore_color.rgb = PowerPointStyles.LIGHT_BLUE
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = PowerPointStyles.BLACK
                    run.font.size = Pt(7)

            cell = table.cell(i, 2)
            cell.text = median_value
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = PowerPointStyles.BLACK
                    run.font.size = Pt(7)

    def _add_tam_cell(self, slide: Slide, market_analysis: MarketAnalysis, left: Inches, top: Inches, width: Inches, height: Inches) -> None:
        result_12min = next((r for r in market_analysis.car_parc_results if r.drive_time_minutes == 12), None)

        if result_12min:
            tam = int(result_12min.car_parc * result_12min.tam_percentage)
        else:
            tam = 0

        assumed_members = 5000
        implied_washes = tam / assumed_members if tam > 0 else 0
        remaining_monthlies = max(tam - market_analysis.total_market_members, 0)
        remaining_washes = remaining_monthlies / assumed_members if remaining_monthlies > 0 else 0

        tam_data = [
            ("Total Addressable Market", f"{tam:,}"),
            ("Assumed Monthly Members / Car Wash", f"{assumed_members:,}"),
            ("Implied Car Washes Allowed in Market", f"{implied_washes:.1f}"),
            ("Current Monthlies in Market", f"{market_analysis.total_market_members:,}"),
            ("Remaining Monthlies in Market", f"{remaining_monthlies:,}"),
            ("Implied Remaining Car Washes Allowed in Market", f"{remaining_washes:.1f}"),
        ]

        rows = len(tam_data)
        cols = 2
        table = slide.shapes.add_table(rows, cols, left, top, width, Inches(1.0)).table
        self._reset_table_styles(table)

        table.columns[0].width = Inches(2.8)
        table.columns[1].width = Inches(0.7)

        for i, (label, value) in enumerate(tam_data):
            is_title = i == 0
            is_implied_washes = i == 2
            is_remaining_washes = i == 5

            cell = table.cell(i, 0)
            cell.text = label

            if is_implied_washes:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PowerPointStyles.LIGHT_BLUE
            elif is_remaining_washes:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PowerPointStyles.DARK_BLUE

            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = is_title or is_implied_washes or is_remaining_washes
                    run.font.size = Pt(7)
                    if is_remaining_washes:
                        run.font.color.rgb = PowerPointStyles.WHITE
                    else:
                        run.font.color.rgb = PowerPointStyles.BLACK

            cell = table.cell(i, 1)
            cell.text = value

            if is_implied_washes:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PowerPointStyles.LIGHT_BLUE
            elif is_remaining_washes:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PowerPointStyles.DARK_BLUE

            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(7)
                    run.font.bold = is_title or is_implied_washes or is_remaining_washes
                    if is_remaining_washes:
                        run.font.color.rgb = PowerPointStyles.WHITE
                    else:
                        run.font.color.rgb = PowerPointStyles.BLACK

    def _add_satellite_cell(self, slide: Slide, market_analysis: MarketAnalysis, left: Inches, top: Inches, width: Inches, height: Inches) -> None:
        if self._location_service:
            coordinates = Coordinates(
                latitude=market_analysis.latitude,
                longitude=market_analysis.longitude,
                formatted_address=market_analysis.address,
            )
            satellite_bytes = self._location_service.download_satellite_image(coordinates, self._satellite_zoom_level)

            if satellite_bytes:
                satellite_stream = BytesIO(satellite_bytes)
                slide.shapes.add_picture(satellite_stream, left, top, width=width, height=height)
            else:
                logger.warning("Failed to download satellite image")
        else:
            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            text_frame.text = "No image available"
            text_frame.paragraphs[0].font.size = Pt(10)

    def _add_site_score_cell(self, slide: Slide, market_analysis: MarketAnalysis, left: Inches, top: Inches, width: Inches, height: Inches) -> None:
        car_counts = market_analysis.key_stats.car_counts or ""
        population = market_analysis.key_stats.population or ""
        median_income = market_analysis.key_stats.median_income
        income_str = f"${median_income:,.0f}" if median_income else ""

        criteria_scores = [
            ("Car Counts", 0, 0.075, f"{car_counts:,}" if car_counts else ""),
            ("Visibility", 0, 0.15, ""),
            ("Layout", 0, 0.05, ""),
            ("Ease", 0, 0.05, ""),
            ("Vac Space", 0, 0.05, ""),
            ("Retail", 0, 0.125, ""),
            ("Quality Competition", 0, 0.225, ""),
            ("Population", 0, 0.225, f"{population:,}" if population else ""),
            ("Income", 0, 0.05, income_str),
        ]

        weighted_score = sum(score * weight for _, score, weight, _ in criteria_scores)
        grade = self._lookup_grade(weighted_score)

        criteria_data = [("Site Score", "Score", "Weighting", "Statistic")]
        for label, score, weight, stat in criteria_scores:
            criteria_data.append((label, score, f"{weight * 100:.1f}%", stat))
        criteria_data.append(("Weighted Score", f"{weighted_score:.1f}", "", ""))
        criteria_data.append(("Grade", grade, "", ""))

        rows = len(criteria_data)
        cols = 4
        table = slide.shapes.add_table(rows, cols, left, top + Inches(0.15), Inches(2.8), Inches(2.2)).table
        self._reset_table_styles(table)

        table.columns[0].width = Inches(1.1)
        table.columns[1].width = Inches(0.5)
        table.columns[2].width = Inches(0.6)
        table.columns[3].width = Inches(0.6)

        for i, (label, score, weight, stat) in enumerate(criteria_data):
            is_header = i == 0
            is_summary = i >= len(criteria_data) - 2
            is_criteria_row = 1 <= i <= len(criteria_scores)

            cell = table.cell(i, 0)
            cell.text = label
            if is_summary:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PowerPointStyles.LIGHT_BLUE
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = PowerPointStyles.BLACK
                    run.font.bold = is_header or is_summary
                    run.font.size = Pt(8)

            cell = table.cell(i, 1)
            cell.text = str(score) if not is_header else score
            if is_criteria_row and isinstance(score, (int, float)):
                self._apply_score_color_scale(cell, score)
            elif is_summary:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PowerPointStyles.LIGHT_BLUE
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = PowerPointStyles.BLACK
                    run.font.bold = is_header or is_summary
                    run.font.size = Pt(8)

            cell = table.cell(i, 2)
            cell.text = weight
            if is_summary:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PowerPointStyles.LIGHT_BLUE
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = PowerPointStyles.BLACK
                    run.font.bold = is_header or is_summary
                    run.font.size = Pt(8)

            cell = table.cell(i, 3)
            cell.text = stat
            if is_summary:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PowerPointStyles.LIGHT_BLUE
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = PowerPointStyles.BLACK
                    run.font.bold = is_header
                    run.font.size = Pt(8)

    def _apply_score_color_scale(self, cell: _Cell, score: float) -> None:
        clamped = max(5, min(10, score))
        value = (clamped - 5) / 5

        start = (0xFF, 0x5D, 0x5D)
        end = (0x00, 0xB0, 0x50)

        r = int(start[0] + (end[0] - start[0]) * value)
        g = int(start[1] + (end[1] - start[1]) * value)
        b = int(start[2] + (end[2] - start[2]) * value)

        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(r, g, b)

    def _lookup_grade(self, weighted_score: float) -> str:
        grading_scale = [
            (9.5, "A+"),
            (9.0, "A"),
            (8.5, "A-"),
            (8.0, "B+"),
            (7.5, "B"),
            (7.0, "B-"),
            (6.5, "C+"),
            (6.0, "C"),
            (5.5, "C-"),
            (5.0, "D"),
            (0.0, "F"),
        ]
        for threshold, grade in grading_scale:
            if weighted_score >= threshold:
                return grade
        return "F"

    def _add_drive_time_cell(self, slide: Slide, market_analysis: MarketAnalysis, left: Inches, top: Inches, width: Inches, height: Inches) -> None:
        priority_times = [10, 12, 15]
        priority_results = [r for r in market_analysis.car_parc_results if r.drive_time_minutes in priority_times]
        priority_results.sort(key=lambda x: priority_times.index(x.drive_time_minutes))

        if priority_results:
            rows = 6
            cols = len(priority_results) + 1
            table = slide.shapes.add_table(rows, cols, left, top, Inches(4.5), Inches(1.2)).table
            self._reset_table_styles(table)

            table.columns[0].width = Inches(1.8)
            for col_idx in range(1, cols):
                table.columns[col_idx].width = Inches(0.9)

            headers = ["", "10 mins", "12 mins", "15 mins"]
            for col_idx, header in enumerate(headers[:cols]):
                cell = table.cell(0, col_idx)
                cell.text = header
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = PowerPointStyles.BLACK
                        run.font.bold = True
                        run.font.size = Pt(8)

            row_labels = [
                "Total Car Parc",
                "Total Addressable Members",
                "Total Addressable Market",
                "Washville Market Share",
                "Washville Site Level Monthlies",
            ]

            for row_idx, label in enumerate(row_labels, start=1):
                table.cell(row_idx, 0).text = label
                for col_idx, result in enumerate(priority_results, start=1):
                    if row_idx == 1:
                        table.cell(row_idx, col_idx).text = f"{result.car_parc:,}"
                    elif row_idx == 2:
                        table.cell(row_idx, col_idx).text = f"{result.tam_percentage * 100:.0f}%"
                    elif row_idx == 3:
                        tam = int(result.car_parc * result.tam_percentage)
                        table.cell(row_idx, col_idx).text = f"{tam:,}"
                    elif row_idx == 4:
                        table.cell(row_idx, col_idx).text = f"{result.market_share_percentage * 100:.0f}%"
                    elif row_idx == 5:
                        monthlies = int(result.car_parc * result.tam_percentage * result.market_share_percentage)
                        table.cell(row_idx, col_idx).text = f"{monthlies:,}"

                for c in range(cols):
                    cell = table.cell(row_idx, c)
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(8)

    def _add_budget_cell(self, slide: Slide, market_analysis: MarketAnalysis, left: Inches, top: Inches, width: Inches, height: Inches) -> None:
        budget_data = [
            ("Preliminary Budget", ""),
            ("Land Cost", ""),
            ("Build Cost", ""),
            ("Equipment/Install/DRB Cost", ""),
            ("Total Cost", ""),
            ("Expected EBITDA (Year 3)", ""),
            ("Unlevered Cash-on-Cash Return (Target > 25%)", ""),
        ]

        rows = len(budget_data)
        cols = 2
        table = slide.shapes.add_table(rows, cols, left, top + Inches(0.15), Inches(2.3), Inches(1.3)).table
        self._reset_table_styles(table)

        table.columns[0].width = Inches(1.8)
        table.columns[1].width = Inches(0.5)

        for i, (label, value) in enumerate(budget_data):
            is_header = i == 0
            is_highlighted = i >= 4

            cell = table.cell(i, 0)
            cell.text = label
            if is_highlighted:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PowerPointStyles.LIGHT_BLUE
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = PowerPointStyles.BLACK
                    run.font.bold = is_header or is_highlighted
                    run.font.size = Pt(8)

            cell = table.cell(i, 1)
            cell.text = value
            if is_highlighted:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PowerPointStyles.LIGHT_BLUE
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = PowerPointStyles.BLACK
                    run.font.bold = is_header or is_highlighted
                    run.font.size = Pt(8)

    def _add_logo_to_slide(self, slide: Slide) -> None:
        if not LOGO_PATH.exists():
            logger.warning(f"Logo file not found at {LOGO_PATH}")
            return

        try:
            left = Inches(8.3)
            top = Inches(0.2)
            width = Inches(1.5)

            slide.shapes.add_picture(str(LOGO_PATH), left, top, width=width)
        except Exception as e:
            logger.error(f"Failed to add logo to slide: {e}")

    def _reset_table_styles(self, table: Table) -> None:
        for row in table.rows:
            row.height = Inches(0.2)
            for cell in row.cells:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PowerPointStyles.WHITE
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.color.rgb = PowerPointStyles.BLACK
                    paragraph.line_spacing = 1.0
                    paragraph.space_before = Pt(0)
                    paragraph.space_after = Pt(0)

                cell.margin_top = Inches(0.01)
                cell.margin_bottom = Inches(0.01)
                cell.margin_left = Inches(0.03)
                cell.margin_right = Inches(0.03)

                tc = cell._tc
                tc_pr = tc.get_or_add_tcPr()
                for border_name in ["lnL", "lnR", "lnT", "lnB"]:
                    border = tc_pr.find(qn(f"a:{border_name}"))
                    if border is not None:
                        tc_pr.remove(border)
                    new_border = etree.SubElement(tc_pr, qn(f"a:{border_name}"))
                    new_border.set("w", "0")
                    new_border.set("cap", "flat")
                    new_border.set("cmpd", "sng")
                    new_border.set("algn", "ctr")
                    etree.SubElement(new_border, qn("a:noFill"))

    def _apply_color_scale(self, cell: _Cell, value: Optional[float]) -> None:
        if value is None:
            return

        value = max(0, min(1, value))

        if value <= 0.5:
            ratio = value / 0.5
            start = (0xF8, 0x69, 0x6B)
            mid = (0xFF, 0xEB, 0x84)
            r = int(start[0] + (mid[0] - start[0]) * ratio)
            g = int(start[1] + (mid[1] - start[1]) * ratio)
            b = int(start[2] + (mid[2] - start[2]) * ratio)
        else:
            ratio = (value - 0.5) / 0.5
            mid = (0xFF, 0xEB, 0x84)
            end = (0x63, 0xBE, 0x7B)
            r = int(mid[0] + (end[0] - mid[0]) * ratio)
            g = int(mid[1] + (end[1] - mid[1]) * ratio)
            b = int(mid[2] + (end[2] - mid[2]) * ratio)

        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(r, g, b)

    def _get_short_address(self, full_address: str) -> str:
        address_parts = full_address.split(",")
        if len(address_parts) >= 2:
            return address_parts[1].strip()
        return full_address
